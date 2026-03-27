#!/usr/bin/env python3
"""Generate user manual PPTX files for Lunch Pair app."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
ORANGE = RGBColor(0xF9, 0x73, 0x16)
ORANGE_LIGHT = RGBColor(0xFF, 0xF7, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x37, 0x41, 0x51)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)


def add_bg(slide, color=ORANGE_LIGHT):
    """Add background color to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=0):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=12, color=BLACK):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def create_title_slide(prs, title, subtitle):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, WHITE)

    # Orange header bar
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(3.5), ORANGE)

    # App icon
    add_text_box(slide, Inches(3.5), Inches(0.6), Inches(3), Inches(1.2),
                 "🍽️", font_size=60, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(8), Inches(1),
                 title, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(2.6), Inches(8), Inches(0.6),
                 subtitle, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Version info
    add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
                 "Version 1.01 | 2026.03", font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


def create_section_slide(prs, number, title, emoji=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Orange accent bar
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(5.63), ORANGE)

    # Section number
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(1),
                 f"{number:02d}", font_size=60, color=ORANGE, bold=True)

    # Title with emoji
    display = f"{emoji}  {title}" if emoji else title
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(8), Inches(1),
                 display, font_size=32, color=BLACK, bold=True)


def create_content_slide(prs, title, bullets, note=""):
    """Create a content slide with title and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), ORANGE)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
                 title, font_size=22, color=WHITE, bold=True)

    # Bullet content
    add_bullet_list(slide, Inches(0.7), Inches(1.2), Inches(8.5), Inches(3.8),
                    bullets, font_size=14)

    # Note at bottom
    if note:
        add_shape(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(0.8), ORANGE_LIGHT)
        add_text_box(slide, Inches(0.7), Inches(4.7), Inches(8.6), Inches(0.6),
                     f"💡 {note}", font_size=11, color=DARK)

    return slide


def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Create a two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), ORANGE)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
                 title, font_size=22, color=WHITE, bold=True)

    # Left column
    add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.3), Inches(0.5), ORANGE_LIGHT)
    add_text_box(slide, Inches(0.6), Inches(1.25), Inches(4), Inches(0.4),
                 left_title, font_size=14, color=ORANGE, bold=True)
    add_bullet_list(slide, Inches(0.6), Inches(1.8), Inches(4), Inches(3),
                    left_items, font_size=12)

    # Right column
    add_shape(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.5), ORANGE_LIGHT)
    add_text_box(slide, Inches(5.4), Inches(1.25), Inches(4), Inches(0.4),
                 right_title, font_size=14, color=ORANGE, bold=True)
    add_bullet_list(slide, Inches(5.4), Inches(1.8), Inches(4), Inches(3),
                    right_items, font_size=12)


def create_table_slide(prs, title, headers, rows):
    """Create a table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), ORANGE)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
                 title, font_size=22, color=WHITE, bold=True)

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * num_rows)
    )
    table = table_shape.table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)

    return slide


# =============================================================
# JAPANESE VERSION
# =============================================================
def generate_ja():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    # --- Slide 1: Title ---
    create_title_slide(prs,
        "Lunch Pair ユーザーマニュアル",
        "社内ランチマッチングアプリ 操作ガイド")

    # --- Slide 2: TOC ---
    create_content_slide(prs, "目次", [
        "01. Lunch Pair とは",
        "02. アカウント登録・ログイン",
        "03. プロフィール設定",
        "04. スワイプ（相手を探す）",
        "05. マッチ＆トーク",
        "06. いいね確認（プレミアム機能）",
        "07. マイページ（プロフィール編集）",
        "08. 通報機能",
        "09. よくある質問（FAQ）",
    ])

    # --- Slide 3: Section 1 ---
    create_section_slide(prs, 1, "Lunch Pair とは", "🍽️")

    # --- Slide 4: What is ---
    create_content_slide(prs, "Lunch Pair とは", [
        "社内の同僚とランチパートナーをマッチングするアプリです",
        "ティンダーのようなスワイプ操作で、気軽に相手を選べます",
        "お互いが「いいかも！」を選んだらマッチ成立",
        "マッチすると相手の顔写真・名前が公開され、チャットが可能に",
        "ランチの好みや曜日・時間・予算で絞り込みもできます",
        "",
        "🔒 プライバシー保護",
        "  • マッチ前：写真のみ表示（名前は非公開）",
        "  • マッチ後：顔写真・名前・詳細情報が公開",
    ], note="本アプリは社内POCとして運用しています。")

    # --- Slide 5: Section 2 ---
    create_section_slide(prs, 2, "アカウント登録・ログイン", "🔐")

    # --- Slide 6: Registration ---
    create_content_slide(prs, "アカウント登録", [
        "1. アプリにアクセスすると、ログイン画面が表示されます",
        "2.「アカウントを作成」をタップ",
        "3. 会社のメールアドレスとパスワードを入力",
        "4. 登録したメールアドレスに認証コードが届きます",
        "5. 認証コードを入力して登録完了",
        "",
        "⚠️ パスワードは8文字以上で設定してください",
        "⚠️ 認証コードの有効期限は24時間です",
    ])

    # --- Slide 7: Section 3 ---
    create_section_slide(prs, 3, "プロフィール設定", "👤")

    # --- Slide 8: Profile Setup ---
    create_two_column_slide(prs, "プロフィール設定（初回ログイン時）",
        "基本情報（必須）", [
            "✅ 表示名（ニックネーム）",
            "✅ 写真1（メイン写真）",
            "✅ 写真2（サブ写真）",
            "📷 写真3・4（任意）",
            "🏢 所属部署（任意）",
        ],
        "ランチ設定（任意）", [
            "📅 ランチ可能曜日（月〜金）",
            "🕐 ランチ時間（11:30〜13:00）",
            "💰 予算（~500円〜~2000円）",
            "📍 エリア（本社周辺・駅周辺 等）",
        ])

    # --- Slide 9: Preferences ---
    create_content_slide(prs, "こだわり設定", [
        "以下の食の好みを複数選択できます：",
        "",
        "🍣 和食　　🍝 洋食　　🥟 中華　　🌶️ 韓国料理",
        "🍛 カレー　🍜 ラーメン　🍣 寿司　☕ カフェ",
        "🥗 ヘルシー　🍖 がっつり　🔥 辛いもの好き　☕ コーヒー好き",
        "",
        "自由記入欄もあります（例：「ベジタリアンです」「辛さ控えめ希望」）",
        "",
        "⭐ 共通のこだわりが多い相手ほど、カードが上位に表示されます",
    ])

    # --- Slide 10: Section 4 ---
    create_section_slide(prs, 4, "スワイプ（相手を探す）", "👆")

    # --- Slide 11: Swipe ---
    create_two_column_slide(prs, "スワイプ操作",
        "操作方法", [
            "→ 右スワイプ＝「いいかも！」（OK）",
            "← 左スワイプ＝「また今度」（SKIP）",
            "カードには写真・部署・共通こだわりが表示",
            "フィルターボタンで条件を絞り込める",
            "",
            "💡 共通こだわりが多い順に表示",
        ],
        "スワイプ制限", [
            "🆓 無料会員：1日3回まで",
            "⭐ プレミアム：無制限",
            "",
            "制限到達時のメッセージ：",
            "「本日のスワイプ上限に達しました」",
            "",
            "💡 毎日0時にリセットされます",
        ])

    # --- Slide 12: Section 5 ---
    create_section_slide(prs, 5, "マッチ＆トーク", "💬")

    # --- Slide 13: Match ---
    create_content_slide(prs, "マッチの仕組み", [
        "お互いが「いいかも！」を選んだ場合にマッチが成立します",
        "",
        "🎉 マッチ成立時：",
        "  • 「マッチしました！」のお祝い画面が表示",
        "  • 相手の顔写真と名前が見れるようになります",
        "  • トーク画面でチャットが可能に",
        "",
        "📱 マッチ一覧（トークタブ）：",
        "  • マッチした相手の一覧が表示されます",
        "  • 最新のメッセージがプレビュー表示されます",
        "  • 未読メッセージがある場合、赤いバッジで件数表示",
    ])

    # --- Slide 14: Chat ---
    create_content_slide(prs, "トーク（チャット）機能", [
        "マッチした相手と1対1でメッセージのやり取りができます",
        "",
        "📝 メッセージ送信：テキスト入力して送信ボタンをタップ",
        "📋 会話のきっかけ：初回は提案メッセージが用意されています",
        "   • 「今週のランチ、一緒にいかがですか？」",
        "   • 「おすすめのお店はありますか？」",
        "   • 「〇〇が好きなんですね！私もです！」等",
        "",
        "👤 相手のプロフィール：ヘッダーの名前をタップで詳細表示",
        "🚩 通報：右上の旗アイコンから通報可能",
    ], note="メッセージはリアルタイムで更新されます（約1秒間隔）")

    # --- Slide 15: Section 6 ---
    create_section_slide(prs, 6, "いいね確認", "💛")

    # --- Slide 16: Likes ---
    create_two_column_slide(prs, "いいね確認（お相手からタブ）",
        "お相手から（受信したいいね）", [
            "🔒 無料会員：ロック状態",
            "⭐ プレミアム：解除して閲覧可能",
            "",
            "解除すると：",
            "• あなたにOKした人の一覧が表示",
            "• 「いいかも」or「スキップ」で返答",
            "• いいかもを返すとマッチ成立！",
        ],
        "自分から（送ったいいね）", [
            "全員閲覧可能（無料機能）",
            "",
            "表示内容：",
            "• 🟢 マッチ済み：名前・写真表示",
            "• ⏳ 返事待ち：名前非公開",
            "",
            "💡 マッチ済みの相手はトークで会話可能",
        ])

    # --- Slide 17: Section 7 ---
    create_section_slide(prs, 7, "マイページ", "⚙️")

    # --- Slide 18: My Page ---
    create_content_slide(prs, "マイページ（プロフィール編集）", [
        "マイページでは以下の操作ができます：",
        "",
        "📝 プロフィール編集",
        "  • 表示名の変更",
        "  • 写真の追加・変更・削除",
        "  • こだわり（食の好み）の変更",
        "  • ランチ設定（曜日・時間・予算・エリア）の変更",
        "  • 倫理的価値観の設定",
        "",
        "🔒 アカウント操作",
        "  • ログアウト",
    ])

    # --- Slide 19: Section 8 ---
    create_section_slide(prs, 8, "通報機能", "🚩")

    # --- Slide 20: Report ---
    create_content_slide(prs, "通報機能", [
        "不適切な行為を発見した場合、ユーザーを通報できます",
        "",
        "通報理由（選択式）：",
        "  1. 不適切なメッセージ",
        "  2. なりすまし・偽プロフィール",
        "  3. 不快な写真",
        "  4. スパム・宣伝行為",
        "  5. ハラスメント・嫌がらせ",
        "  6. 個人情報の悪用",
        "  7. ドタキャン・無断キャンセル",
        "  8. その他",
        "",
        "💡 詳細を記入いただくと、対応がスムーズです",
    ], note="通報内容は運営チームが確認し、適切に対応します。")

    # --- Slide 21: Section 9 ---
    create_section_slide(prs, 9, "よくある質問（FAQ）", "❓")

    # --- Slide 22: FAQ ---
    create_content_slide(prs, "よくある質問", [
        "Q. 1日にスワイプできる回数は？",
        "A. 無料会員は1日3回まで。毎日0時にリセットされます。",
        "",
        "Q. マッチする前に相手の名前は分かる？",
        "A. いいえ。マッチが成立するまで名前は非公開です。",
        "",
        "Q. メッセージはリアルタイム？",
        "A. はい。約1秒間隔で自動更新されます。",
        "",
        "Q. パスワードを忘れた場合は？",
        "A. ログイン画面の「パスワードを忘れた場合」から再設定できます。",
        "",
        "Q. 問い合わせはどこから？",
        "A. マイページから問い合わせフォームをご利用ください。",
    ])

    # --- Slide 23: End ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.63), ORANGE)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(1),
                 "🍽️", font_size=60, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(0.8),
                 "Lunch Pair", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.3), Inches(8), Inches(0.5),
                 "素敵なランチパートナーを見つけましょう！",
                 font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.3), Inches(8), Inches(0.5),
                 "お問い合わせ：マイページ → お問い合わせ",
                 font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    prs.save("docs/user-manual-ja.pptx")
    print("✅ docs/user-manual-ja.pptx generated (23 slides)")


# =============================================================
# KOREAN VERSION
# =============================================================
def generate_ko():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)

    # --- Slide 1: Title ---
    create_title_slide(prs,
        "Lunch Pair 사용자 매뉴얼",
        "사내 점심 매칭 앱 사용 가이드")

    # --- Slide 2: TOC ---
    create_content_slide(prs, "목차", [
        "01. Lunch Pair란?",
        "02. 계정 등록 및 로그인",
        "03. 프로필 설정",
        "04. 스와이프 (상대 찾기)",
        "05. 매칭 & 채팅",
        "06. 좋아요 확인 (프리미엄 기능)",
        "07. 마이페이지 (프로필 편집)",
        "08. 신고 기능",
        "09. 자주 묻는 질문 (FAQ)",
    ])

    # --- Slide 3: Section 1 ---
    create_section_slide(prs, 1, "Lunch Pair란?", "🍽️")

    # --- Slide 4: What is ---
    create_content_slide(prs, "Lunch Pair란?", [
        "사내 동료와 점심 파트너를 매칭해주는 앱입니다",
        "틴더 스타일의 스와이프 조작으로 간편하게 상대를 선택합니다",
        "서로 '좋아요!'를 선택하면 매칭 성립",
        "매칭되면 상대의 사진·이름이 공개되고 채팅이 가능해집니다",
        "점심 취향, 요일, 시간, 예산으로 필터링도 가능합니다",
        "",
        "🔒 프라이버시 보호",
        "  • 매칭 전: 사진만 표시 (이름 비공개)",
        "  • 매칭 후: 사진·이름·상세 정보 공개",
    ], note="본 앱은 사내 POC로 운영 중입니다.")

    # --- Slide 5: Section 2 ---
    create_section_slide(prs, 2, "계정 등록 및 로그인", "🔐")

    # --- Slide 6: Registration ---
    create_content_slide(prs, "계정 등록", [
        "1. 앱에 접속하면 로그인 화면이 표시됩니다",
        "2. '계정 만들기'를 탭",
        "3. 회사 이메일 주소와 비밀번호를 입력",
        "4. 등록한 이메일로 인증 코드가 발송됩니다",
        "5. 인증 코드를 입력하면 등록 완료",
        "",
        "⚠️ 비밀번호는 8자 이상으로 설정해주세요",
        "⚠️ 인증 코드의 유효기간은 24시간입니다",
    ])

    # --- Slide 7: Section 3 ---
    create_section_slide(prs, 3, "프로필 설정", "👤")

    # --- Slide 8: Profile ---
    create_two_column_slide(prs, "프로필 설정 (최초 로그인 시)",
        "기본 정보 (필수)", [
            "✅ 표시 이름 (닉네임)",
            "✅ 사진 1 (메인 사진)",
            "✅ 사진 2 (서브 사진)",
            "📷 사진 3·4 (선택)",
            "🏢 소속 부서 (선택)",
        ],
        "점심 설정 (선택)", [
            "📅 점심 가능 요일 (월~금)",
            "🕐 점심 시간 (11:30~13:00)",
            "💰 예산 (~500엔 ~ ~2000엔)",
            "📍 에리어 (본사 주변·역 주변 등)",
        ])

    # --- Slide 9: Preferences ---
    create_content_slide(prs, "취향 설정 (こだわり)", [
        "다음 음식 취향을 복수 선택할 수 있습니다:",
        "",
        "🍣 일식　　🍝 양식　　🥟 중식　　🌶️ 한식",
        "🍛 카레　🍜 라멘　🍣 스시　☕ 카페",
        "🥗 건강식　🍖 든든한 식사　🔥 매운 음식 좋아함　☕ 커피 좋아함",
        "",
        "자유 기입란도 있습니다 (예: '비건입니다', '매운 거 약하게 희망')",
        "",
        "⭐ 공통 취향이 많은 상대일수록 카드가 상위에 표시됩니다",
    ])

    # --- Slide 10: Section 4 ---
    create_section_slide(prs, 4, "스와이프 (상대 찾기)", "👆")

    # --- Slide 11: Swipe ---
    create_two_column_slide(prs, "스와이프 조작",
        "조작 방법", [
            "→ 오른쪽 스와이프 = '좋아요!' (OK)",
            "← 왼쪽 스와이프 = '다음에' (SKIP)",
            "카드에는 사진·부서·공통 취향이 표시",
            "필터 버튼으로 조건 필터링 가능",
            "",
            "💡 공통 취향이 많은 순서로 표시",
        ],
        "스와이프 제한", [
            "🆓 무료 회원: 하루 3회까지",
            "⭐ 프리미엄: 무제한",
            "",
            "제한 도달 시 메시지:",
            "'오늘의 스와이프 상한에 도달했습니다'",
            "",
            "💡 매일 자정에 리셋됩니다",
        ])

    # --- Slide 12: Section 5 ---
    create_section_slide(prs, 5, "매칭 & 채팅", "💬")

    # --- Slide 13: Match ---
    create_content_slide(prs, "매칭 구조", [
        "서로 '좋아요!'를 선택한 경우 매칭이 성립됩니다",
        "",
        "🎉 매칭 성립 시:",
        "  • '매칭되었습니다!' 축하 화면 표시",
        "  • 상대의 사진과 이름이 보이게 됩니다",
        "  • 토크 화면에서 채팅 가능",
        "",
        "📱 매칭 목록 (토크 탭):",
        "  • 매칭된 상대 목록이 표시됩니다",
        "  • 최신 메시지가 미리보기로 표시됩니다",
        "  • 안 읽은 메시지가 있으면 빨간 뱃지로 건수 표시",
    ])

    # --- Slide 14: Chat ---
    create_content_slide(prs, "채팅 기능", [
        "매칭된 상대와 1대1로 메시지를 주고받을 수 있습니다",
        "",
        "📝 메시지 전송: 텍스트 입력 후 전송 버튼 탭",
        "📋 대화 시작 도우미: 첫 대화 시 제안 메시지가 준비되어 있습니다",
        "   • '이번 주 점심, 같이 어떠세요?'",
        "   • '추천 식당 있으세요?'",
        "   • '〇〇 좋아하시는군요! 저도요!' 등",
        "",
        "👤 상대 프로필: 헤더의 이름을 탭하면 상세 표시",
        "🚩 신고: 오른쪽 상단 깃발 아이콘에서 신고 가능",
    ], note="메시지는 실시간으로 업데이트됩니다 (약 1초 간격)")

    # --- Slide 15: Section 6 ---
    create_section_slide(prs, 6, "좋아요 확인", "💛")

    # --- Slide 16: Likes ---
    create_two_column_slide(prs, "좋아요 확인 (お相手から 탭)",
        "받은 좋아요", [
            "🔒 무료 회원: 잠금 상태",
            "⭐ 프리미엄: 해제하여 열람 가능",
            "",
            "해제하면:",
            "• 나에게 OK한 사람 목록 표시",
            "• '좋아요' 또는 '스킵'으로 응답",
            "• 좋아요를 반환하면 매칭 성립!",
        ],
        "보낸 좋아요", [
            "전원 열람 가능 (무료 기능)",
            "",
            "표시 내용:",
            "• 🟢 매칭 완료: 이름·사진 표시",
            "• ⏳ 답장 대기: 이름 비공개",
            "",
            "💡 매칭 완료 상대는 채팅으로 대화 가능",
        ])

    # --- Slide 17: Section 7 ---
    create_section_slide(prs, 7, "마이페이지", "⚙️")

    # --- Slide 18: My Page ---
    create_content_slide(prs, "마이페이지 (프로필 편집)", [
        "마이페이지에서는 다음 작업이 가능합니다:",
        "",
        "📝 프로필 편집",
        "  • 표시 이름 변경",
        "  • 사진 추가·변경·삭제",
        "  • 취향 (음식 선호) 변경",
        "  • 점심 설정 (요일·시간·예산·에리어) 변경",
        "  • 윤리적 가치관 설정",
        "",
        "🔒 계정 조작",
        "  • 로그아웃",
    ])

    # --- Slide 19: Section 8 ---
    create_section_slide(prs, 8, "신고 기능", "🚩")

    # --- Slide 20: Report ---
    create_content_slide(prs, "신고 기능", [
        "부적절한 행위를 발견한 경우 사용자를 신고할 수 있습니다",
        "",
        "신고 사유 (선택식):",
        "  1. 부적절한 메시지",
        "  2. 사칭·가짜 프로필",
        "  3. 불쾌한 사진",
        "  4. 스팸·광고 행위",
        "  5. 괴롭힘·따돌림",
        "  6. 개인정보 악용",
        "  7. 약속 무단 취소",
        "  8. 기타",
        "",
        "💡 상세 내용을 기입하면 대응이 빨라집니다",
    ], note="신고 내용은 운영팀이 확인하여 적절히 대응합니다.")

    # --- Slide 21: Section 9 ---
    create_section_slide(prs, 9, "자주 묻는 질문 (FAQ)", "❓")

    # --- Slide 22: FAQ ---
    create_content_slide(prs, "자주 묻는 질문", [
        "Q. 하루에 스와이프할 수 있는 횟수는?",
        "A. 무료 회원은 하루 3회까지. 매일 자정에 리셋됩니다.",
        "",
        "Q. 매칭 전에 상대의 이름을 알 수 있나요?",
        "A. 아니요. 매칭이 성립될 때까지 이름은 비공개입니다.",
        "",
        "Q. 메시지는 실시간인가요?",
        "A. 네. 약 1초 간격으로 자동 업데이트됩니다.",
        "",
        "Q. 비밀번호를 잊었을 경우?",
        "A. 로그인 화면의 '비밀번호를 잊었을 경우'에서 재설정 가능합니다.",
        "",
        "Q. 문의는 어디서?",
        "A. 마이페이지에서 문의 폼을 이용해주세요.",
    ])

    # --- Slide 23: End ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.63), ORANGE)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(1),
                 "🍽️", font_size=60, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(0.8),
                 "Lunch Pair", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.3), Inches(8), Inches(0.5),
                 "멋진 점심 파트너를 찾아보세요!",
                 font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.3), Inches(8), Inches(0.5),
                 "문의: 마이페이지 → 문의하기",
                 font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    prs.save("docs/user-manual-ko.pptx")
    print("✅ docs/user-manual-ko.pptx generated (23 slides)")


if __name__ == "__main__":
    generate_ja()
    generate_ko()
